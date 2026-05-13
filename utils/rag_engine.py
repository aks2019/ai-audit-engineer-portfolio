import pandas as pd
import json
from io import StringIO
from datetime import datetime
import os
import hashlib
from pathlib import Path
from langchain_postgres import PGVector
from dotenv import load_dotenv
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import SystemMessage
from langchain_google_genai import ChatGoogleGenerativeAI
os.environ["LANGCHAIN_TRACING_V2"] = "false"       # comment out while using online AI Model
os.environ["LANGCHAIN_API_KEY"] = "dummy"           # prevents any accidental trace, comment out while using online AI Model
from langchain_openai import ChatOpenAI             # comment out while using online AI Model
from langchain_community.chat_models import ChatOpenAI
from langchain_core.documents import Document
from utils.redis_cache import cache_rag_result
load_dotenv()

_embeddings = None
_embeddings_error = None


def _get_embeddings():
    """Lazy-init embeddings to avoid app startup delay/network errors."""
    global _embeddings
    global _embeddings_error
    if _embeddings_error is not None:
        raise RuntimeError(_embeddings_error)
    if _embeddings is None:
        try:
            from langchain_huggingface import HuggingFaceEmbeddings
            # Default local-only avoids external SSL/download retry loops.
            # Set HF_LOCAL_FILES_ONLY=false if internet + cert chain is configured.
            local_only = os.getenv("HF_LOCAL_FILES_ONLY", "true").lower() == "true"
            _embeddings = HuggingFaceEmbeddings(
                model_name="all-MiniLM-L6-v2",
                model_kwargs={"local_files_only": local_only},
            )
        except Exception as e:
            _embeddings_error = (
                "HuggingFace embeddings unavailable. "
                "Set HF_LOCAL_FILES_ONLY=false with internet/cert access, "
                "or pre-download the model locally. "
                f"Original error: {e}"
            )
            raise RuntimeError(_embeddings_error) from e
    return _embeddings

def get_vectorstore():
    return PGVector(
        embeddings=_get_embeddings(),
        collection_name="audit_policies",
        connection=os.getenv("NEON_CONNECTION_STRING"),
        use_jsonb=True
    )

def _get_vectorstore_safe():
    """Returns (vectorstore, error_message). error_message is None on success."""
    try:
        return get_vectorstore(), None
    except Exception as e:
        return None, str(e)

def _get_llm():
    # If USE_LOCAL_LLM is explicitly set to "true", use local llama.cpp; otherwise use Gemini cloud
    if os.getenv("USE_LOCAL_LLM", "true").lower() == "true":
        from langchain_openai import ChatOpenAI
        return ChatOpenAI(base_url="http://127.0.0.1:8080/v1", api_key="llama.cpp", # https://integrate.api.nvidia.com/v1 
                          model="gemma-4-E2B-it-Q4_K_M.gguf", temperature=0.7)      # minimaxai/minimax-m2.7  
    # Otherwise, use Gemini (Cloud LLM)
    from langchain_google_genai import ChatGoogleGenerativeAI
    return ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.7,
                                  google_api_key=os.getenv("GOOGLE_API_KEY"))

def _get_db_path():
    """Get path to local SQLite audit database."""
    Path("data").mkdir(exist_ok=True)
    return "data/audit.db"

# ====================== STANDARDS REGISTRY SEARCH (SQLite) ======================

def search_standards_from_sqlite(query: str, top_k: int = 8) -> list:
    """
    Search the audit_standards SQLite table by keyword match on reference,
    description, applicability, and clause_text.
    
    Returns list of dicts with keys: family, reference, description, 
    applicability, clause_text, source_url
    
    This is a fast, local, zero-cost alternative to PGVector for structured
    audit standards.
    """
    import sqlite3
    try:
        conn = sqlite3.connect(_get_db_path())
        cursor = conn.cursor()
        
        # Split query into keywords — fix: group word+number combos (e.g. "AS 2" → ["AS 2"])
        # Also strip punctuation from tokens (e.g. "AS 13," becomes "AS 13")
        raw_tokens = query.replace("'", "").replace('"', "").split()
        keywords = []
        i = 0
        while i < len(raw_tokens):
            # Strip punctuation from current token
            curr = raw_tokens[i].strip(",;:.?!")
            # Strip punctuation from next token if exists
            nxt = raw_tokens[i+1].strip(",;:.?!") if i + 1 < len(raw_tokens) else ""
            
            # If current token is alpha and next token is a number, combine them (e.g. "AS" "2" → "AS 2")
            # Also handle cases like "13," → cleaned to "13" which IS a digit
            if (nxt and curr.isalpha() and nxt.isdigit()):
                keywords.append(f"{curr} {nxt}")
                i += 2
            else:
                keywords.append(curr)
                i += 1
        
        # Build query: match any keyword across multiple columns
        conditions = []
        params = []
        for kw in keywords:
            if len(kw) < 2:
                continue
            like = f"%{kw}%"
            conditions.append(
                "(reference LIKE ? OR description LIKE ? OR clause_text LIKE ? OR applicability LIKE ?)"
            )
            params.extend([like, like, like, like])
        
        if not conditions:
            conn.close()
            return []
        
        # Build a smart ORDER BY that:
        # 1. Gives highest priority to keyword(s) that look like a reference (e.g. "AS 2", "Section 188")
        # 2. Then by how many different keywords matched the row
        # 3. Then alphabetically by family/reference
        order_cases = []
        order_params = []
        ref_like_kws = []
        for kw in keywords:
            # Detect "reference-like" keywords (e.g. "AS 2", "Section 188", "Ind AS 116", "Clause 12")
            is_ref_like = any(
                kw.upper().startswith(prefix) or kw.upper() == prefix.rstrip()
                for prefix in ["AS", "IND AS", "SA", "SIA", "CAS", "CARO", "SECTION", "CLAUSE"]
            )
            if is_ref_like or any(ch.isdigit() for ch in kw):
                ref_like_kws.append(kw)
        
        # Priority 1: ALL exact reference matches first (e.g. reference = 'AS 2', reference = 'AS 13')
        for rk in ref_like_kws:
            order_cases.append("CASE WHEN reference = ? THEN 0 ELSE 1 END")
            order_params.append(rk)
        
        # Priority 2: ALL LIKE reference matches (e.g. reference LIKE '%Clause 12%')
        for rk in ref_like_kws:
            order_cases.append("CASE WHEN reference LIKE ? THEN 0 ELSE 1 END")
            order_params.append(f"%{rk}%")
        
        # Priority 3: Number of matched keywords (higher = more relevant)
        n_kw = len(keywords)
        for j in range(n_kw):
            order_cases.append(f"CASE WHEN (reference LIKE ? OR description LIKE ? OR clause_text LIKE ? OR applicability LIKE ?) THEN 0 ELSE 1 END")
            kw_j = keywords[j] if j < len(keywords) else ""
            order_params.extend([f"%{kw_j}%"] * 4)
        
        # Priority 3: Alphabetical as tiebreaker
        order_str = ", ".join(order_cases) if order_cases else "1"
        order_str += ", family, reference"
        
        sql = f"""
            SELECT family, reference, description, applicability, clause_text, source_url
            FROM audit_standards
            WHERE {' OR '.join(conditions)}
            ORDER BY {order_str}
            LIMIT ?
        """
        params = params + order_params + [top_k]
        
        cursor.execute(sql, params)
        rows = cursor.fetchall()
        conn.close()
        
        results = []
        for row in rows:
            results.append({
                "family": row[0],
                "reference": row[1],
                "description": row[2],
                "applicability": row[3],
                "clause_text": row[4],
                "source_url": row[5],
            })
        return results
    except Exception as e:
        return [{"family": "ERROR", "reference": str(e), "description": "Database error", 
                 "applicability": "", "clause_text": "", "source_url": ""}]

def format_standards_for_context(standards: list) -> str:
    """Format matched standards into a compact string for LLM context."""
    if not standards:
        return ""
    
    lines = ["─── AUDIT STANDARDS REGISTRY (Matches from Companies Act / CARO / Ind AS / AS / CAS / SIA / SA) ───"]
    for i, std in enumerate(standards, 1):
        clause = std.get("clause_text", "") or std.get("description", "")
        lines.append(
            f"  [{i}] {std['family']} — {std['reference']}"
            f"\n      Requirement: {clause[:200]}"
            f"\n      Applies to: {std.get('applicability', 'All companies')}"
        )
    lines.append("─── END OF STANDARDS REGISTRY ───\n")
    return "\n".join(lines)

def search_standards_by_family(family: str) -> list:
    """Get all standards for a given family (e.g. 'Companies Act', 'Ind AS', 'CARO')."""
    import sqlite3
    try:
        conn = sqlite3.connect(_get_db_path())
        cursor = conn.cursor()
        cursor.execute(
            "SELECT family, reference, description, applicability, clause_text, source_url "
            "FROM audit_standards WHERE family = ? ORDER BY reference",
            (family,)
        )
        rows = cursor.fetchall()
        conn.close()
        return [
            {
                "family": r[0], "reference": r[1], "description": r[2],
                "applicability": r[3], "clause_text": r[4], "source_url": r[5]
            }
            for r in rows
        ]
    except Exception:
        return []

# ====================== COMPLIANCE KEYWORD EXTRACTOR ======================

def extract_compliance_keywords_from_data(flagged_transactions: list) -> str:
    """
    Extract compliance-relevant keywords from flagged transaction data
    for searching the Standards Registry.
    
    This is MUCH more effective than searching using the LLM prompt text
    (which contains noise words like 'generate', 'summary', 'executive').
    """
    keywords = set()
    
    # Always include core audit compliance areas
    core_areas = [
        "payment", "vendor", "overdue", "credit", "terms",
        "related party", "board", "approval", "disclosure",
        "internal control", "risk", "audit", "fraud",
        "financial", "reporting", "compliance", "threshold",
        "asset", "liability", "revenue", "expense",
        "contract", "procurement", "purchase", "invoice",
        "reconciliation", "verification", "documentation"
    ]
    for c in core_areas:
        keywords.add(c)
    
    # Extract from transaction data
    for txn in flagged_transactions[:30]:
        # Related party flag
        for rp_field in ["related_party", "related_party_flag", "is_related_party", "rpt"]:
            val = str(txn.get(rp_field, "") or "")
            if val.lower() in ["1", "true", "yes", "y"]:
                keywords.add("related party")
                keywords.add("section 188")
                keywords.add("CARO clause 12")
        
        # Overdue detection
        for od_field in ["days_overdue", "overdue_days", "overdue", "ageing_days", "aging_days"]:
            overdue = txn.get(od_field, 0) or 0
            if isinstance(overdue, (int, float)) and overdue > 30:
                keywords.add("overdue")
                keywords.add("payment terms")
                keywords.add("credit period")
                if overdue > 180:
                    keywords.add("doubtful debt")
                    keywords.add("impairment")
                elif overdue > 90:
                    keywords.add("going concern")
        
        # Amount thresholds
        amount = txn.get("amount", 0) or txn.get("payment_amount", 0) or 0
        if isinstance(amount, (int, float)):
            if amount > 10000000:  # > 1 Cr
                keywords.add("board approval")
                keywords.add("section 186")
            if amount > 5000000:  # > 50 L
                keywords.add("materiality")
                keywords.add("disclosure")
        
        # Category/area
        for cat_field in ["category", "area", "type", "expense_type", "description"]:
            val = str(txn.get(cat_field, "") or "")
            if val and val.lower() not in ["unknown", "none", "", "nan"]:
                clean = val.lower().replace("_", " ").replace("-", " ")
                for w in clean.split():
                    if len(w) > 2:
                        keywords.add(w)
        
        # Vendor name
        vendor = str(txn.get("vendor_name", "") or txn.get("vendor", "") or "")
        # Skip synthetic names like VENDOR_001
        if vendor and not any(pat in vendor.upper() for pat in ["VENDOR_", "VEND_", "V_", "TEST"]):
            keywords.add(vendor.lower().replace("_", " "))
        
        # High anomaly score
        for score_field in ["anomaly_score", "anomaly_probability", "xgb_risk_score", "risk_score"]:
            score = txn.get(score_field, 0) or 0
            if isinstance(score, (int, float)) and score > 0.8:
                keywords.add("anomaly detection")
                keywords.add("audit evidence")
                keywords.add("SA 240")  # fraud standard
    
    return " ".join(sorted(keywords))


# ====================== UPDATED RAG PROMPT (now standards-aware) ======================

@cache_rag_result(ttl_seconds=1800)
def get_rag_chain():
    llm = _get_llm()
    RAG_PROMPT = ChatPromptTemplate.from_messages([
        SystemMessage(content="""You are AI AUDIT BOT – a Continuous Control Monitoring + Policy Compliance Agent for Emami Agrotech Limited.
You are a senior internal auditor with 17+ years FMCG experience.

You have TWO knowledge sources available in the context BELOW (separated by section headers):

**Source 1: AUDIT STANDARDS REGISTRY** (look for the section starting with "─── AUDIT STANDARDS REGISTRY ───")
- These contain the actual regulatory requirements from Companies Act, CARO, Ind AS, AS, CAS, SIA, SA
- Example: "Companies Act — Section 188: Prior Board/SR approval for related party transactions"
- **You MUST reference these standards by name in your answer** (e.g. "Section 188 of Companies Act 2013 requires...")

**Source 2: UPLOADED POLICY DOCUMENTS** (look for the section starting with "─── UPLOADED POLICY DOCUMENTS ───")
- These are company's internal policies, SOPs, contracts

**YOUR JOB:**
1. Read the flagged transaction data in the question
2. Check the Standards Registry section — identify which standards apply to the risk areas found
3. Check the Policy Documents section — identify which internal rules apply
4. Compare the flagged data against BOTH to determine compliance gaps
5. Cite EXACT standard references (Companies Act Section X, CARO Clause X, Ind AS X, SA X, etc.)

STRICT RESPONSE FORMAT (never deviate):
1. Direct Answer
2. Root Cause Analysis (if applicable; otherwise "Not applicable as this is a direct policy query.")
3. Risk/Compliance Implication
4. Policy/Contract Reference(s)
   * Document name + Page + exact clause/table
   * Standard registry reference (e.g. Companies Act Section 188, CARO Clause 12)
5. Recommended Next Audit Step

Rules:
- Answer using the provided context. Never hallucinate or add information.
- **If you see a Standards Registry section in the context, you MUST reference it.**
- Keep every section concise and audit-ready.
- Always end with "📚 Sources & References"
- Use bullet points only inside sections.
- Never add extra headings or explanations."""),
        ("human", """Context (Standards Registry + Policy Documents):
{context}

Question: {question}
Answer:""")
    ])
    return RAG_PROMPT | llm

@cache_rag_result(ttl_seconds=1800)
def get_free_form_chain():
    llm = _get_llm()
    FREE_PROMPT = ChatPromptTemplate.from_messages([
        SystemMessage(content="""You are AI AUDIT BOT – a senior internal auditor with 17+ years FMCG experience.

You have TWO knowledge sources available in the context:
1. **Audit Standards Registry** (Companies Act, CARO, Ind AS, AS, CAS, SIA, SA) — the LAW/REGULATIONS
2. **Uploaded Policy Documents** (company policies, SOPs, contracts) — COMPANY INTERNAL RULES

**IMPORTANT: You MUST search the Standards Registry section in the context and reference any applicable standards by name (Section number, Clause number, Ind AS number, etc.)**

Always verify uploaded data against relevant standards. Cite exact references.
Always end with "📚 Sources & References" if documents were used."""),
        ("human", """Context (Standards Registry + Policy Documents):
{context}

Question: {question}
Answer:""")
    ])
    return FREE_PROMPT | llm

# ====================== DOCUMENT UPLOAD (unchanged) ======================

def add_documents_from_upload(uploaded_file):  # PDF only
    from langchain_community.document_loaders import PyMuPDFLoader
    import tempfile
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp_path = tmp.name
    loader = PyMuPDFLoader(tmp_path)
    docs = loader.load()
    get_vectorstore().add_documents(docs)
    os.unlink(tmp_path)
    return len(docs)

def add_documents_from_csv_or_excel_or_office(uploaded_file):  # CSV, XLSX, DOCX, PPTX
    import tempfile
    suffix = Path(uploaded_file.name).suffix.lower()
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.getvalue())
            tmp_path = tmp.name

        if suffix in [".csv", ".xlsx"]:
            import pandas as pd
            df = pd.read_csv(tmp_path) if suffix == ".csv" else pd.read_excel(tmp_path)
            docs = [Document(page_content=row.to_string(), metadata={"source": uploaded_file.name, "row": i}) 
                    for i, row in df.iterrows()]
        elif suffix == ".docx":
            import docx2txt
            text = docx2txt.process(tmp_path)
            docs = [Document(page_content=text, metadata={"source": uploaded_file.name})]
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            docs = [Document(page_content=text, metadata={"source": uploaded_file.name})]
        else:
            docs = []

        get_vectorstore().add_documents(docs)
        return len(docs)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

# ====================== COMBINED SEARCH (PGVector + SQLite Standards) ======================

def build_combined_context(query: str, pgvector_top_k: int = 4, standards_top_k: int = 8,
                          flagged_data: list = None) -> tuple:
    """
    Search BOTH PGVector (policy docs) and SQLite (standards registry).
    
    For standards search: uses flagged_data keywords if available (extracts
    compliance-relevant terms like 'related party', 'overdue', 'board approval'
    from the transaction data). Falls back to query text if no flagged_data.
    
    For PGVector search: uses the original query which works better for semantic
    similarity (PGVector understands full sentences).
    
    Returns:
        (combined_context_string, citations_list)
    
    The LLM receives: Standards Registry matches + Policy Document matches
    """
    citations = []
    
    # ── 1. Search SQLite Standards Registry ──────────────────────────
    # Use flagged transaction data to extract REAL compliance keywords
    if flagged_data:
        standards_search_query = extract_compliance_keywords_from_data(flagged_data)
    else:
        standards_search_query = query  # fallback to original query
    
    matched_standards = search_standards_from_sqlite(standards_search_query, top_k=standards_top_k)
    standards_context = format_standards_for_context(matched_standards)
    
    if matched_standards and matched_standards[0]["family"] != "ERROR":
        for std in matched_standards:
            citations.append(f"Standards Registry: {std['family']} — {std['reference']}")
    
    # ── 2. Search PGVector Policy Documents (uses original query for semantic search) ──
    vectorstore, db_error = _get_vectorstore_safe()
    pgvector_context = ""
    pgvector_docs = []
    
    if vectorstore is not None:
        try:
            pgvector_docs = vectorstore.similarity_search(query[:1000], k=pgvector_top_k)
            pgvector_context = "\n\n".join(
                f"Document: {d.metadata.get('source', 'Policy doc')}\n{d.page_content}"
                for d in pgvector_docs
            )
            for doc in pgvector_docs:
                source = doc.metadata.get("source", "pgvector policy/contract")
                page = doc.metadata.get("page", "")
                citations.append(f"{source} (page {page})" if page else f"{source}")
        except Exception:
            pgvector_context = "(Policy documents unavailable — PGVector query failed)"
    else:
        pgvector_context = "(Policy documents unavailable — PGVector connection error)"
    
    # ── 3. Combine Both ──────────────────────────────────────────────
    combined = f"""{standards_context}

─── UPLOADED POLICY DOCUMENTS (Company Internal Rules / SOPs / Contracts) ───
{pgvector_context}
─── END OF POLICY DOCUMENTS ───
"""
    return combined, citations


# ====================== UPDATED: Generate RAG Audit Report (now standards-aware) ======================

def generate_rag_audit_report(flagged_transactions, contract_text=None, vendor_name=None):
    # Cap payload to top 15 rows
    sample = flagged_transactions[:15] if isinstance(flagged_transactions, list) else flagged_transactions

    query = f"""
Generate full audit-ready executive summary for these flagged high-risk vendor payments:
{sample}

Vendor: {vendor_name or 'Unknown'}
Include specific clause violations from any attached contract or policy documents.
Focus on related-party, payment terms, ageing, financial limits, and anomaly thresholds.

For each finding, reference:
- The specific regulatory standard violated (Companies Act Section, CARO Clause, Ind AS, etc.)
- The internal policy/contract clause that was breached
"""
    if contract_text:
        query += f"\n\nATTACHED VENDOR CONTRACT TEXT (analyse for violations):\n{contract_text[:4000]}"
    
    # ── Build combined context ──────────────────────────────────────
    # PASS the flagged transactions data for keyword extraction
    combined_context, citations = build_combined_context(
        query, pgvector_top_k=4, standards_top_k=10,
        flagged_data=flagged_transactions if isinstance(flagged_transactions, list) else None
    )
    
    # ── If both sources are empty/unavailable ─────────────────────────
    if not combined_context.strip():
        return {
            "audit_summary": (
                "⚠️ **No context available** — could not reach PGVector or SQLite.\n\n"
                "Check your Neon connection and ensure data/audit.db exists locally."
            ),
            "citations": [],
            "log_hash": "offline",
        }

    # ── LLM Call ────────────────────────────────────────────────────
    try:
        chain = get_rag_chain()
        response = chain.invoke({"context": combined_context, "question": query})
    except Exception as e:
        return {
            "audit_summary": f"⚠️ **LLM call failed** — local model did not respond.\n\n`{e}`\n\nEnsure llama.cpp is running on `http://127.0.0.1:8080` and the model is loaded.",
            "citations": citations,
            "log_hash": "llm-error",
        }

    audit_summary = response.content if hasattr(response, "content") else str(response)
    log_hash = hashlib.sha256(audit_summary.encode()).hexdigest()[:16]

    return {
        "audit_summary": audit_summary,
        "citations": citations,
        "log_hash": log_hash
    }


# ====================== ROBUST LLM-POWERED FINANCIAL STATEMENT DRAFTER ======================

def load_tb_raw(tb_file) -> pd.DataFrame:
    """Simple load + net calculation only."""
    if tb_file.name.endswith('.xlsx'):
        df = pd.read_excel(tb_file)
    else:
        df = pd.read_csv(tb_file)
    df['period_net'] = df['period_credit'].fillna(0) - df['period_debit'].fillna(0)
    df['closing_net'] = df['closing_credit'].fillna(0) - df['closing_debit'].fillna(0)
    return df

def llm_draft_financial_statements(df: pd.DataFrame, company_name: str = "Emami Agrotech Limited") -> dict:
    """LLM drafts full statements with Manufacturing Account + Trading + P&L + BS + approx Cash Flow."""
    summary_df = df.groupby('account_group').agg({
        'period_net': 'sum',
        'closing_net': 'sum',
        'account_name': 'count'
    }).round(2).reset_index()
    summary_text = summary_df.to_string(index=False)
    
    prompt = f"""You are a senior CA with 17+ years FMCG internal audit experience at Emami Agrotech.
Given the SAP Trial Balance grouped summary below, draft the **complete Indian manufacturing financial statements** in STRICT JSON format only.

TB Summary:
{summary_text}

Required sections (use standard FMCG format):
- Manufacturing Account (cost of production)
- Trading Account
- Profit & Loss Account
- Balance Sheet
- Cash Flow Statement (approximate from TB changes)

Return ONLY valid JSON. No extra text.
{{
  "manufacturing": {{"Raw Materials Consumed": number, "Direct Expenses": number, "Manufacturing Overheads": number, "Cost of Production": number}},
  "trading": {{"Sales": number, "Opening Stock": number, "COGS / Purchases": number, "Direct Expenses": number, "Closing Stock": number, "Gross Profit": number}},
  "pl": {{"Gross Profit": number, "Other Income": number, "Indirect Expenses": number, "Depreciation": number, "Net Profit": number}},
  "bs": {{"Fixed Assets (Net)": number, "Current Assets": number, "Total Assets": number, "Current Liabilities": number, "Borrowings": number, "Equity & Reserves": number, "Total Liabilities & Equity": number}},
  "cash_flow_approx": {{"Operating Activities": number, "Investing Activities": number, "Financing Activities": number, "Net Cash Flow": number}},
  "major_heads": [array of objects from summary]
}}"""

    chain = get_rag_chain()
    response = chain.invoke({"context": "SAP Trial Balance data", "question": prompt})
    
    # Robust JSON extraction
    raw_output = response.content if hasattr(response, "content") else str(response)
    try:
        json_str = raw_output.strip()
        if "```json" in json_str:
            json_str = json_str.split("```json")[1].split("```")[0].strip()
        elif "```" in json_str:
            json_str = json_str.split("```")[1].strip()
        result = json.loads(json_str)
    except Exception:
        result = {"error": "JSON parse failed", "raw_llm_output": raw_output[:2000]}
    
    return result

def generate_financial_audit_report(drafted: dict, company_name: str = "Emami Agrotech Limited"):
    query = f"""
    Audit these LLM-drafted financial statements for {company_name}:
    Manufacturing: {drafted.get('manufacturing', {})}
    Trading: {drafted.get('trading', {})}
    P&L: {drafted.get('pl', {})}
    BS: {drafted.get('bs', {})}
    Cash Flow (approx): {drafted.get('cash_flow_approx', {})}
    Major heads: {drafted.get('major_heads', [])}
    
    For each financial statement area, reference relevant:
    - Ind AS or AS standards (e.g. Ind AS 1 Presentation, Ind AS 16 PPE, Ind AS 116 Leases)
    - Companies Act sections (Sections 128-134)
    - CARO clauses if applicable
    """
    
    # Build combined context
    combined_context, citations = build_combined_context(
        query, pgvector_top_k=8, standards_top_k=10
    )
    
    chain = get_rag_chain()
    response = chain.invoke({"context": combined_context, "question": query})
    audit_summary = response.content if hasattr(response, "content") else str(response)
    
    log_hash = hashlib.sha256(audit_summary.encode()).hexdigest()[:16]
    
    return {
        "audit_summary": audit_summary,
        "citations": citations,
        "log_hash": log_hash,
        **drafted
    }