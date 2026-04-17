import streamlit as st
import json
import os
from datetime import datetime
from pathlib import Path   # ← ADD THIS LINE
from utils.rag_engine import get_vectorstore, get_free_form_chain, get_rag_chain, add_documents_from_upload, add_documents_from_csv_or_excel_or_office, _get_vectorstore_safe
from db_utils import log_rag_query
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from io import BytesIO
import io as io_module
import tempfile

st.title("📋 AI Policy RAG Bot")
st.caption("Continuous Control Monitoring + Policy Compliance Agent | 100% Audit Trail | Built by Ashok Sharma")

SAVED_CHATS_DIR = "saved_chats"
os.makedirs(SAVED_CHATS_DIR, exist_ok=True)
CURRENT_CHAT_FILE = "current_chat.json"

def load_current_chat():
    if os.path.exists(CURRENT_CHAT_FILE):
        with open(CURRENT_CHAT_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return []

def save_current_chat(messages):
    with open(CURRENT_CHAT_FILE, "w", encoding="utf-8") as f:
        json.dump(messages, f, ensure_ascii=False, indent=2)

if "messages" not in st.session_state:
    st.session_state.messages = load_current_chat()
if "chat_title" not in st.session_state:
    st.session_state.chat_title = f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"

# Sidebar (all original controls preserved)
with st.sidebar:
    st.header("Chat Controls")
    if st.button("➕ New Conversation"):
        if st.session_state.messages:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"chat_{timestamp}.json"
            path = os.path.join(SAVED_CHATS_DIR, filename)
            with open(path, "w", encoding="utf-8") as f:
                json.dump(st.session_state.messages, f, ensure_ascii=False, indent=2)
            st.success(f"Old chat auto-saved as {filename}")
        st.session_state.messages = []
        st.session_state.chat_title = f"Chat {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        save_current_chat(st.session_state.messages)
        st.rerun()

    if st.button("🗑️ Clear Current Chat"):
        st.session_state.messages = []
        save_current_chat(st.session_state.messages)
        st.rerun()

    if st.button("📥 Export Current Chat to PDF"):
        pdf_buffer = generate_pdf(st.session_state.messages, st.session_state.chat_title)
        st.download_button("Download PDF", pdf_buffer, f"{st.session_state.chat_title}.pdf", "application/pdf")

    saved_files = sorted([f for f in os.listdir(SAVED_CHATS_DIR) if f.endswith(".json")], reverse=True)
    if saved_files:
        st.subheader("Saved Conversations")
        selected = st.selectbox("Load saved chat", ["None"] + saved_files)
        if selected != "None" and st.button("Load"):
            path = os.path.join(SAVED_CHATS_DIR, selected)
            with open(path, "r", encoding="utf-8") as f:
                st.session_state.messages = json.load(f)
            st.session_state.chat_title = selected.replace(".json", "")
            save_current_chat(st.session_state.messages)
            st.rerun()

    st.subheader("Search Chat History")
    search_term = st.text_input("Search questions & answers", "")
    if search_term:
        matches = []
        for msg in st.session_state.messages + load_current_chat():
            if search_term.lower() in msg["content"].lower():
                role = "User" if msg["role"] == "user" else "Bot"
                matches.append(f"**{role}**: {msg['content'][:120]}...")
        if matches:
            st.write("Found matches:")
            for m in matches[:10]:
                st.markdown(m)
        else:
            st.info("No matches found")

def generate_pdf(messages, title="Conversation"):
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    y = height - 50
    c.setFont("Helvetica-Bold", 14)
    c.drawString(50, y, f"AI Policy RAG Bot - {title}")
    y -= 30
    c.setFont("Helvetica", 10)
    c.drawString(50, y, f"Exported: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    y -= 40
    for msg in messages:
        role = "User" if msg["role"] == "user" else "Bot"
        c.setFont("Helvetica-Bold", 11)
        c.drawString(50, y, f"{role}:")
        y -= 20
        c.setFont("Helvetica", 10)
        text = msg["content"].replace("\n", " ").strip()
        text_lines = [text[i:i+90] for i in range(0, len(text), 90)]
        for line in text_lines:
            if y < 50:
                c.showPage()
                y = height - 50
            c.drawString(70, y, line)
            y -= 15
        y -= 20
    c.save()
    buffer.seek(0)
    return buffer


def _extract_chat_attachment_text(uploaded_files):
    """Read text from user chat attachments only (does not index the vector store)."""
    if not uploaded_files:
        return ""
    parts = []
    for f in uploaded_files:
        name = f.name
        suffix = Path(name).suffix.lower()
        try:
            if suffix == ".pdf":
                from langchain_community.document_loaders import PyMuPDFLoader

                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(f.getvalue())
                    tmp_path = tmp.name
                try:
                    loader = PyMuPDFLoader(tmp_path)
                    docs = loader.load()
                    text = "\n".join(d.page_content for d in docs)
                finally:
                    os.unlink(tmp_path)
            elif suffix == ".txt":
                text = f.getvalue().decode("utf-8", errors="replace")
            elif suffix == ".csv":
                import pandas as pd

                df = pd.read_csv(io_module.BytesIO(f.getvalue()))
                text = df.to_string()
            elif suffix == ".xlsx":
                import pandas as pd

                df = pd.read_excel(io_module.BytesIO(f.getvalue()))
                text = df.to_string()
            elif suffix == ".docx":
                import docx2txt

                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(f.getvalue())
                    tmp_path = tmp.name
                try:
                    text = docx2txt.process(tmp_path)
                finally:
                    os.unlink(tmp_path)
            elif suffix == ".pptx":
                from pptx import Presentation

                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tmp.write(f.getvalue())
                    tmp_path = tmp.name
                try:
                    prs = Presentation(tmp_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                finally:
                    os.unlink(tmp_path)
            else:
                text = f"(Unsupported file type: {name})"
        except Exception as ex:
            text = f"(Could not read {name}: {ex})"
        parts.append(f"--- File: {name} ---\n{text}")
    return "\n\n".join(parts)


st.markdown(f"**Active Chat:** {st.session_state.chat_title}")
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# UPLOAD SUPPORT FOR PDF + CSV + EXCEL + DOCX + PPTX
uploaded_doc = st.file_uploader(
    "Upload new policy/contract PDF / SAP CSV / Excel / DOCX / PPTX",
    type=["pdf", "csv", "xlsx", "docx", "pptx"]
)
if uploaded_doc:
    with st.spinner("Indexing..."):
        suffix = Path(uploaded_doc.name).suffix.lower()
        if suffix == ".pdf":
            count = add_documents_from_upload(uploaded_doc)
        else:
            count = add_documents_from_csv_or_excel_or_office(uploaded_doc)
        st.success(f"✅ Indexed {count} pages/rows into pgvector")
# Response Mode radio (preserved)
mode = st.radio("Response Mode", ["Structured Audit Report", "Free-Form Discussion"], horizontal=True, key="response_mode")

if "chat_attachment_widget_key" not in st.session_state:
    st.session_state.chat_attachment_widget_key = 0

attach_col, input_col = st.columns([1, 11], vertical_alignment="bottom", gap="small")
attachment_files = None
with attach_col:
    # Gemini / ChatGPT–style: one control (Material +) beside the composer; file picker opens in the panel.
    with st.popover(
        "",
        icon=":material/add:",
        help="Attach files for this message only. They are not added to the policy index.",
        use_container_width=True,
    ):
        st.caption("Drop files here or choose from your device")
        attachment_files = st.file_uploader(
            "chat_attach_pick",
            type=["pdf", "csv", "xlsx", "docx", "pptx", "txt"],
            accept_multiple_files=True,
            key=f"chat_attach_{st.session_state.chat_attachment_widget_key}",
            label_visibility="collapsed",
        )
        if attachment_files:
            st.caption("Selected")
            for uf in attachment_files:
                st.caption(f"• {uf.name}")
with input_col:
    prompt = st.chat_input("Ask about policy, clause, contract...")

if prompt:
    attachment_text = _extract_chat_attachment_text(attachment_files) if attachment_files else ""
    user_content = prompt
    if attachment_files:
        names = ", ".join(f.name for f in attachment_files)
        user_content = f"{prompt}\n\n*Attachments: {names}*"

    st.session_state.messages.append({"role": "user", "content": user_content})
    with st.chat_message("user"):
        st.markdown(user_content)

    with st.chat_message("assistant"):
        with st.spinner("Searching..."):
            # === FIXED: Inject full conversation history ===
            history = "\n".join([f"{m['role'].capitalize()}: {m['content']}" for m in st.session_state.messages[:-1]])
            
            vectorstore, db_error = _get_vectorstore_safe()
            if vectorstore is None:
                st.error(f"⚠️ Policy database unavailable: {db_error}")
                st.stop()
            docs = vectorstore.similarity_search(prompt, k=5)
            context = "\n\n---\n\n".join(f"Document: {d.metadata.get('source', 'Policy')}\n{d.page_content}" for d in docs)
            
            attachment_block = ""
            if attachment_text:
                attachment_block = f"""User-uploaded attachments (this message only; not from the indexed policy store):
{attachment_text}

"""

            full_context = f"""{attachment_block}Previous conversation history:
{history}

Retrieved policy/contract documents:
{context}"""
            
            chain = get_rag_chain() if mode == "Structured Audit Report" else get_free_form_chain()
            response = chain.invoke({"context": full_context, "question": prompt})
            
            st.markdown(response.content)
            with st.expander("📚 Sources & References"):
                for i, doc in enumerate(docs, 1):
                    st.markdown(f"**{i}.** {doc.metadata.get('source', 'Policy doc')}")
            log_rag_query(prompt, response.content)

    st.session_state.messages.append({"role": "assistant", "content": response.content})
    save_current_chat(st.session_state.messages)
    st.session_state.chat_attachment_widget_key += 1